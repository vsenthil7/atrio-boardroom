"""
Build the ATRIO Boardroom pitch deck for Milan AI Week 2026 submission.

12 slides, 16:9 widescreen, ATRIO brand (ink #0a0a0a / paper #ffffff / orange #f59e0b accent).
Editorial-style typography matching the SPA front-end.

Inputs read from project docs (sources of truth):
  - docs/AT-Hack0021_Claude_ATRIO_BRDv1_20260518.md
  - docs/AT-Hack0021_Claude_ATRIO_UseCaseCatalogue_20260518.md
  - docs/DEMO_RUNBOOK.md
  - docs/ATRIO_Traceability_LIVE.md

Output (relative to atrio/ repo root):
  - submission_media/atrio-pitch-deck-{stamp}.pptx
  - submission_media/atrio-pitch-deck-{stamp}.pdf
  - submission_media/_backup/{same files}

This script is idempotent. Re-run to regenerate after content edits.
"""
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ---------- Brand ----------
INK = RGBColor(0x0A, 0x0A, 0x0A)          # near-black background
PAPER = RGBColor(0xFF, 0xFF, 0xFF)        # text on dark / panel bg
TEXT_PRIMARY = RGBColor(0xE5, 0xE5, 0xE5) # body text on dark
TEXT_SECONDARY = RGBColor(0x9C, 0xA3, 0xAF) # muted
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)       # ATRIO accent (header strip)
BLUE = RGBColor(0x3B, 0x82, 0xF6)         # scene-step accent
GREEN = RGBColor(0x10, 0xB9, 0x81)        # success
RED = RGBColor(0xDC, 0x26, 0x26)          # blocker / danger
PANEL_BG = RGBColor(0x1A, 0x1A, 0x1A)     # darker panel against INK


# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def fill_slide(slide, color):
    """Solid-fill the slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    bold=False,
    italic=False,
    color=PAPER,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Inter",
):
    """Add a textbox with a single paragraph (multi-line via \\n)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_top_strip(slide, hackathon="AT-HACK0021 · MILAN AI WEEK 2026", color=ORANGE):
    """Editorial-style top strip: orange tick + hackathon label."""
    add_rect(slide, Inches(0.5), Inches(0.35), Inches(0.25), Inches(0.06), color)
    add_text(
        slide,
        Inches(0.85), Inches(0.25), Inches(10), Inches(0.3),
        hackathon,
        size=10, color=TEXT_SECONDARY, font="Inter",
    )


def add_footer(slide, page_num=None, total=12):
    add_text(
        slide,
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
        "ATRIO Boardroom · vsenthil7/atrio-boardroom · Apache 2.0",
        size=9, color=TEXT_SECONDARY, font="Inter",
    )
    if page_num is not None:
        add_text(
            slide,
            Inches(11), Inches(7.05), Inches(2), Inches(0.3),
            f"{page_num} / {total}",
            size=9, color=TEXT_SECONDARY, font="Inter", align=PP_ALIGN.RIGHT,
        )


# ---------- Slide builders ----------


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    # Massive title
    add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(1.5),
             "ATRIO", size=96, bold=True, color=PAPER, font="Inter")
    # Blue underline
    add_rect(s, Inches(0.55), Inches(3.55), Inches(1.5), Inches(0.06), BLUE)
    # Subtitle
    add_text(s, Inches(0.5), Inches(3.85), Inches(12), Inches(1.0),
             "Your AI boardroom.",
             size=44, italic=True, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(4.85), Inches(12), Inches(1.0),
             "Six specialist agents. One audited table. Mandate-enforced treasury.",
             size=22, color=TEXT_PRIMARY, font="Inter")
    add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
             "Milan AI Week 2026 · AT-Hack0021 · github.com/vsenthil7/atrio-boardroom",
             size=12, color=TEXT_SECONDARY, font="Inter")
    return s


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The problem", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Founders and family offices decide alone.",
             size=44, bold=True, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(2.6), Inches(12), Inches(2.5),
             "Big decisions go one of two ways today:\n\n"
             "  → Delegated to a single advisor — fast, but a single point of failure.\n"
             "  → Convened with a committee — slow, hard to schedule, hard to audit.\n\n"
             "No middle option. No audit trail of the reasoning. No way to enforce a\n"
             "mandate at machine speed. No way to replay a decision in six months.",
             size=18, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 2)
    return s


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "The solution", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "An AI boardroom that holds a real debate.",
             size=40, bold=True, color=PAPER, font="Inter")

    # Three columns of value
    cols = [
        ("DEBATE", "Six specialist agents (CFO, CTO, CMO, COO, Counsel, Facilitator)\nargue with each other before answering. Dissent triggers re-runs."),
        ("ENFORCE", "Per-tenant mandate is checked at the API on every treasury action.\nTwo-party authorisation cannot be bypassed by any client."),
        ("AUDIT", "Append-only audit log captures every turn, vote, model invocation,\nand state transition. Exportable as JSONL + manifest for compliance."),
    ]
    col_w = Inches(4.1)
    gap = Inches(0.1)
    for i, (header, body) in enumerate(cols):
        x = Inches(0.5) + (col_w + gap) * i
        add_rect(s, x, Inches(2.8), col_w, Inches(0.06), BLUE)
        add_text(s, x, Inches(2.95), col_w, Inches(0.4),
                 header, size=14, bold=True, color=ORANGE, font="Inter")
        add_text(s, x, Inches(3.45), col_w, Inches(3),
                 body, size=14, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 3)
    return s


def slide_six_agents(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Six specialists", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Distinct personas. Distinct models. One table.",
             size=32, bold=True, color=PAPER, font="Inter")

    agents = [
        ("CFO",         "Capital + burn",          "Gemini 2.5 Flash"),
        ("CTO",         "Tech feasibility",        "Featherless Qwen2-72B"),
        ("CMO",         "Narrative + GTM",         "Featherless Llama-3.1-70B"),
        ("COO",         "Ops + people",            "Featherless Mistral-Large"),
        ("Counsel",     "Legal + regulatory",      "Gemini 2.5 Pro"),
        ("Facilitator", "Synthesis + dissent",     "Gemini 2.5 Pro"),
    ]
    col_w = Inches(4.0)
    row_h = Inches(0.85)
    for i, (role, focus, model) in enumerate(agents):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col_w * col + Inches(0.05) * col
        y = Inches(2.7) + (row_h + Inches(0.2)) * row
        add_rect(s, x, y, col_w, row_h, PANEL_BG, line_color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.05), col_w - Inches(0.4), Inches(0.4),
                 role, size=18, bold=True, color=PAPER, font="Inter")
        add_text(s, x + Inches(0.2), y + Inches(0.4), col_w - Inches(0.4), Inches(0.3),
                 focus, size=11, color=TEXT_PRIMARY, font="Inter")
        add_text(s, x + Inches(0.2), y + Inches(0.6), col_w - Inches(0.4), Inches(0.3),
                 model, size=10, italic=True, color=TEXT_SECONDARY, font="Inter")

    add_text(s, Inches(0.5), Inches(5.6), Inches(12), Inches(1.2),
             "Model registry is the only path to inference. Fallback chains are configured per agent\n"
             "(Gemini primary, Featherless fallback, mock for offline demo). Every model invocation\n"
             "writes to the audit log with token counts, latency, and request ID.",
             size=12, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 4)
    return s


def slide_treasury(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Treasury · Mandate · Two-party", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "The agents can do, not just advise.",
             size=32, bold=True, color=PAPER, font="Inter")

    # State machine: proposed → first_authorised → fully_authorised → executed
    states = [
        ("proposed", ORANGE),
        ("first_authorised", ORANGE),
        ("fully_authorised", BLUE),
        ("executed", GREEN),
    ]
    box_w = Inches(2.6)
    box_h = Inches(0.7)
    gap = Inches(0.4)
    x0 = Inches(0.5)
    y0 = Inches(2.9)
    for i, (state, color) in enumerate(states):
        x = x0 + (box_w + gap) * i
        add_rect(s, x, y0, box_w, box_h, PANEL_BG, line_color=color)
        add_text(s, x, y0 + Inches(0.18), box_w, Inches(0.4),
                 state, size=14, bold=True, color=PAPER, font="Inter", align=PP_ALIGN.CENTER)
        if i < 3:
            # Arrow connector
            add_text(s, x + box_w - Inches(0.05), y0 + Inches(0.12), gap + Inches(0.2), Inches(0.5),
                     "→", size=24, bold=True, color=TEXT_SECONDARY, font="Inter", align=PP_ALIGN.CENTER)

    # Mandate gates panel
    add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.5),
             "Four mandate gates checked at the API layer (not the UI):",
             size=14, color=PAPER, bold=True, font="Inter")
    gates = [
        ("✓ Permitted instrument list",  "SHV-xStock, IEF-xStock, EURUSD-xStock"),
        ("✓ Permitted side",              "buy / sell"),
        ("✓ Single-instrument max",       "€50,000 per proposal"),
        ("✓ Daily loss limit",            "€25,000 rolling 24h"),
    ]
    for i, (gate, detail) in enumerate(gates):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + Inches(6.2) * col
        y = Inches(4.5) + Inches(0.55) * row
        add_text(s, x, y, Inches(3.5), Inches(0.3),
                 gate, size=13, bold=True, color=GREEN, font="Inter")
        add_text(s, x + Inches(3.5), y, Inches(2.5), Inches(0.3),
                 detail, size=12, color=TEXT_SECONDARY, font="Inter", italic=True)

    # Two-party rule banner
    add_rect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.7), PANEL_BG, line_color=RED)
    add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.7),
             "✕ Two-party authorisation cannot be bypassed.\n"
             "The proposer cannot self-second. Tested — and demonstrated live in the recording.",
             size=12, color=PAPER, font="Inter")
    add_footer(s, 5)
    return s


def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Live demo · 4 stages", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "What you will see in the 2-minute video.",
             size=32, bold=True, color=PAPER, font="Inter")

    stages = [
        ("1", "BOARDROOM DEBATE",
         "Founder asks a Q3 hiring question; PDF burn-plan uploaded; six agents debate; dissent re-runs; consensus + action list."),
        ("2", "TREASURY · MANDATE · TWO-PARTY",
         "Founder proposes SHV-xStock buy. Founder authorises (1/2). Founder TRIES to self-second — REFUSED by the API."),
        ("3", "SECOND HUMAN",
         "CEO signs in on a separate browser, same tenant. Authorises (2/2). Trade executes against Kraken paper. Audit gains 5 rows."),
        ("4", "BOARDPACK · AUDIT EXPORT",
         "Session closed. Boardpack PDF download. Audit page. Export ZIP = JSONL + manifest. Compliance-ready."),
    ]
    for i, (num, title, body) in enumerate(stages):
        y = Inches(2.7) + Inches(1.05) * i
        # Big number
        add_text(s, Inches(0.5), y, Inches(0.8), Inches(1),
                 num, size=44, bold=True, color=ORANGE, font="Inter")
        add_text(s, Inches(1.4), y + Inches(0.1), Inches(11), Inches(0.4),
                 title, size=16, bold=True, color=PAPER, font="Inter")
        add_text(s, Inches(1.4), y + Inches(0.5), Inches(11.4), Inches(0.6),
                 body, size=12, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 6)
    return s


def slide_proof(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Proof, not promises", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Audit-grade comes from tests, not slogans.",
             size=32, bold=True, color=PAPER, font="Inter")

    metrics = [
        ("381 / 381",  "backend tests pass"),
        ("90.68 %",    "line coverage (gate 85 %)"),
        ("15 / 15",    "vitest pass · typecheck clean"),
        ("16 / 20",    "Playwright pass against live stack"),
        ("24 / 24",    "verification-a · hard API assertions"),
        ("14 / 14",    "verification-b · OCR rubric items"),
        ("12",         "real bugs found and fixed"),
        ("23",         "commits on main · all green"),
    ]
    box_w = Inches(3.0)
    box_h = Inches(1.4)
    for i, (big, small) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = Inches(0.5) + (box_w + Inches(0.1)) * col
        y = Inches(2.9) + (box_h + Inches(0.2)) * row
        add_rect(s, x, y, box_w, box_h, PANEL_BG, line_color=GREEN if "pass" in small.lower() or "fixed" in small.lower() or "green" in small.lower() else BLUE)
        add_text(s, x, y + Inches(0.2), box_w, Inches(0.6),
                 big, size=28, bold=True, color=PAPER, font="Inter", align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.85), box_w, Inches(0.45),
                 small, size=11, color=TEXT_SECONDARY, font="Inter", align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
             "All numbers reproducible on the repo. `make up && pytest && pwsh ./demovideo/verification-a/run-verify-a.ps1`",
             size=11, color=TEXT_SECONDARY, font="Consolas", italic=True)
    add_footer(s, 7)
    return s


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Architecture", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Boring, predictable, audit-grade.",
             size=32, bold=True, color=PAPER, font="Inter")

    # 3-tier sketch
    tiers = [
        ("FRONTEND",  "React · Vite · Tailwind",      "Caddy SPA on :8080 (PWA-ready)",       BLUE),
        ("API",       "FastAPI · SQLAlchemy · Alembic", "uvicorn on :8000 (RLS · mandate gates)", ORANGE),
        ("INFRA",     "Postgres 16 + pgvector · MinIO · LiveKit · Mailhog", "All in compose · one ``make up``", GREEN),
    ]
    for i, (label, stack, note, color) in enumerate(tiers):
        y = Inches(2.8) + Inches(1.2) * i
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), PANEL_BG, line_color=color)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.4),
                 label, size=16, bold=True, color=color, font="Inter")
        add_text(s, Inches(3.5), y + Inches(0.1), Inches(9), Inches(0.4),
                 stack, size=14, bold=True, color=PAPER, font="Inter")
        add_text(s, Inches(3.5), y + Inches(0.55), Inches(9), Inches(0.4),
                 note, size=11, color=TEXT_SECONDARY, font="Inter", italic=True)

    add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "Sponsor pools wired: Vultr (deploy) · Gemini (inference) · Featherless (fallback) · Speechmatics (voice) · Kraken (treasury paper)",
             size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 8)
    return s


def slide_sponsors(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Sponsor pools utilised", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Five pools. Real integration. No demoware.",
             size=32, bold=True, color=PAPER, font="Inter")

    sponsors = [
        ("Vultr",         "Compute + deployment",              "VM target for the public demo URL"),
        ("Google Gemini", "Primary inference (CFO + Counsel + Facilitator)", "API key configured · model registry routes calls"),
        ("Featherless",   "Specialist fallback (CTO + CMO + COO)",            "Configured · activates when Gemini is down or rate-limited"),
        ("Speechmatics",  "Live STT + diarisation (voice mode)",              "Sidecar wired · text-only fallback in demo recording"),
        ("Kraken",        "Treasury execution (xStocks)",                     "Paper-mode by default · live mode is one config flip"),
    ]
    for i, (name, role, detail) in enumerate(sponsors):
        y = Inches(2.6) + Inches(0.75) * i
        add_text(s, Inches(0.5), y, Inches(2.5), Inches(0.4),
                 name, size=16, bold=True, color=ORANGE, font="Inter")
        add_text(s, Inches(3.2), y, Inches(4.5), Inches(0.4),
                 role, size=13, color=PAPER, font="Inter")
        add_text(s, Inches(7.8), y, Inches(5.2), Inches(0.4),
                 detail, size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_footer(s, 9)
    return s


def slide_what_else(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Beyond the MVP", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "What ships next.",
             size=32, bold=True, color=PAPER, font="Inter")

    columns = [
        ("v1.1 — 3 weeks", [
            "WCAG 2.1 AA via @axe-core/playwright",
            "Performance gates (p95 < 10s session)",
            "Speechmatics live STT (coupon redeemed)",
            "Featherless live fallback chain",
            "Push backend coverage 90.68 % → 100 %",
        ]),
        ("v1.2 — 8 weeks", [
            "Native mobile app (PWA → Capacitor)",
            "IOTA anchoring on treasury receipts",
            "Vultr Object Storage (from MinIO)",
            "Multi-language UI (EN / IT / ES / FR / DE)",
            "Long-horizon memory (per-agent profile)",
        ]),
    ]
    for i, (header, items) in enumerate(columns):
        x = Inches(0.5) + Inches(6.4) * i
        add_rect(s, x, Inches(2.7), Inches(6.0), Inches(0.05), BLUE)
        add_text(s, x, Inches(2.85), Inches(6.0), Inches(0.4),
                 header, size=15, bold=True, color=PAPER, font="Inter")
        for j, item in enumerate(items):
            add_text(s, x, Inches(3.35) + Inches(0.5) * j, Inches(6.0), Inches(0.4),
                     "·  " + item, size=13, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 10)
    return s


def slide_team_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(0.5),
             "Team + ask", size=12, color=ORANGE, font="Inter", bold=True)
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(1.4),
             "Built by Verixa for Milan AI Week 2026.",
             size=32, bold=True, color=PAPER, font="Inter")

    add_text(s, Inches(0.5), Inches(2.8), Inches(12), Inches(0.4),
             "Founder + lead engineer · architecture, infra, demo, deployment",
             size=14, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(3.25), Inches(12), Inches(0.4),
             "Generated under the AT-Hack0021 sprint with Claude as paired engineer.",
             size=12, color=TEXT_SECONDARY, font="Inter", italic=True)

    add_text(s, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
             "What we are asking for",
             size=16, bold=True, color=ORANGE, font="Inter")
    asks = [
        "Vultr credits to host the live demo URL for the EU founder cohort.",
        "Gemini + Featherless production-tier keys for the v1.1 multilingual rollout.",
        "30 minutes with a Kraken xStocks PM to validate the mandate-violation corpus.",
        "Time on stage at Milan AI Week to demo end-to-end with audience-supplied questions.",
    ]
    for i, ask in enumerate(asks):
        add_text(s, Inches(0.5), Inches(5.0) + Inches(0.4) * i, Inches(12), Inches(0.4),
                 "·  " + ask, size=12, color=TEXT_PRIMARY, font="Inter")
    add_footer(s, 11)
    return s


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s, INK)
    add_top_strip(s)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12), Inches(1.4),
             "ATRIO Boardroom",
             size=64, bold=True, color=PAPER, font="Inter")
    add_rect(s, Inches(0.55), Inches(3.4), Inches(1.5), Inches(0.06), BLUE)
    add_text(s, Inches(0.5), Inches(3.65), Inches(12), Inches(1.0),
             "Mandate-enforced at the API.",
             size=28, italic=True, color=PAPER, font="Inter")
    add_text(s, Inches(0.5), Inches(4.3), Inches(12), Inches(1.0),
             "Two-party authorisation. Cross-tenant isolation. Audit-grade by default.",
             size=18, color=TEXT_PRIMARY, font="Inter")

    add_text(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
             "github.com/vsenthil7/atrio-boardroom",
             size=14, color=ORANGE, font="Consolas")
    add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
             "Apache 2.0 · 381 backend tests · 90.68 % coverage · 24/24 + 14/14 demo-verified",
             size=11, color=TEXT_SECONDARY, font="Inter", italic=True)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "Thank you.",
             size=12, color=TEXT_SECONDARY, font="Inter")
    return s


# ---------- Main ----------


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_solution,
        slide_six_agents,
        slide_treasury,
        slide_demo,
        slide_proof,
        slide_architecture,
        slide_sponsors,
        slide_what_else,
        slide_team_thanks,
        slide_closing,
    ]
    for b in builders:
        b(prs)

    # Output dir: atrio/submission_media/ (script is at atrio/scripts/build_pitch_deck.py)
    out_dir = Path(__file__).resolve().parent.parent / "submission_media"
    out_dir.mkdir(parents=True, exist_ok=True)
    backup_dir = out_dir / "_backup"
    backup_dir.mkdir(exist_ok=True)

    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_pptx = out_dir / f"atrio-pitch-deck-{stamp}.pptx"
    prs.save(out_pptx)
    # Backup copy
    import shutil
    shutil.copy2(out_pptx, backup_dir / out_pptx.name)

    print(f"deck: {out_pptx} ({out_pptx.stat().st_size / 1024:.1f} KB)")
    return out_pptx


if __name__ == "__main__":
    build_deck()
